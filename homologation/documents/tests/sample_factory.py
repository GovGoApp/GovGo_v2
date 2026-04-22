from __future__ import annotations

from pathlib import Path
from typing import Dict, List


SAMPLES_DIR = Path(__file__).resolve().parent / "samples"
COMMON_TOKEN = "MARKITDOWN-LAB"


def _write_text(path: Path, content: str) -> None:
    path.write_text(content, encoding="utf-8")


def _build_text_samples() -> List[Dict[str, str]]:
    cases = [
        {
            "name": "txt-local",
            "filename": "sample_markitdown.txt",
            "expected_contains": f"{COMMON_TOKEN} TXT",
            "content": (
                f"{COMMON_TOKEN} TXT\n"
                "GovGo documents homologation sample.\n"
                "Objeto: fornecimento de alimentos hospitalares.\n"
            ),
        },
        {
            "name": "md-local",
            "filename": "sample_markitdown.md",
            "expected_contains": f"{COMMON_TOKEN} MD",
            "content": (
                f"# {COMMON_TOKEN} MD\n\n"
                "- modulo: documentos\n"
                "- origem: laboratorio local\n"
            ),
        },
        {
            "name": "html-local",
            "filename": "sample_markitdown.html",
            "expected_contains": f"{COMMON_TOKEN} HTML",
            "content": (
                "<!doctype html>\n"
                "<html><body>\n"
                f"<h1>{COMMON_TOKEN} HTML</h1>\n"
                "<p>Documento HTML de teste do GovGo.</p>\n"
                "</body></html>\n"
            ),
        },
        {
            "name": "xml-local",
            "filename": "sample_markitdown.xml",
            "expected_contains": f"{COMMON_TOKEN} XML",
            "content": (
                "<?xml version=\"1.0\" encoding=\"UTF-8\"?>\n"
                "<documento>\n"
                f"  <titulo>{COMMON_TOKEN} XML</titulo>\n"
                "  <objeto>Teste de conversao XML.</objeto>\n"
                "</documento>\n"
            ),
        },
        {
            "name": "json-local",
            "filename": "sample_markitdown.json",
            "expected_contains": COMMON_TOKEN,
            "content": (
                "{\n"
                f"  \"titulo\": \"{COMMON_TOKEN} JSON\",\n"
                "  \"objeto\": \"Teste de conversao JSON\",\n"
                "  \"ano\": 2026\n"
                "}\n"
            ),
        },
        {
            "name": "csv-local",
            "filename": "sample_markitdown.csv",
            "expected_contains": COMMON_TOKEN,
            "content": (
                "campo,valor\n"
                f"titulo,{COMMON_TOKEN} CSV\n"
                "objeto,Teste de conversao CSV\n"
            ),
        },
        {
            "name": "tsv-local",
            "filename": "sample_markitdown.tsv",
            "expected_contains": COMMON_TOKEN,
            "content": (
                "campo\tvalor\n"
                f"titulo\t{COMMON_TOKEN} TSV\n"
                "objeto\tTeste de conversao TSV\n"
            ),
        },
        {
            "name": "yaml-local",
            "filename": "sample_markitdown.yaml",
            "expected_contains": COMMON_TOKEN,
            "content": (
                f"titulo: {COMMON_TOKEN} YAML\n"
                "objeto: Teste de conversao YAML\n"
                "ano: 2026\n"
            ),
        },
    ]
    for case in cases:
        _write_text(SAMPLES_DIR / case["filename"], case["content"])
    return cases


def _build_docx_sample() -> Dict[str, str]:
    from docx import Document

    path = SAMPLES_DIR / "sample_markitdown.docx"
    document = Document()
    document.add_heading(f"{COMMON_TOKEN} DOCX", level=1)
    document.add_paragraph("Documento DOCX de teste do laboratorio GovGo.")
    document.add_paragraph("Objeto: aquisicao de insumos hospitalares.")
    document.save(path)
    return {
        "name": "docx-local",
        "filename": path.name,
        "expected_contains": f"{COMMON_TOKEN} DOCX",
        "expected_absent": ["[Content_Types].xml", "<w:document", "styles.xml"],
    }


def _build_pptx_sample() -> Dict[str, str]:
    from pptx import Presentation

    path = SAMPLES_DIR / "sample_markitdown.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = f"{COMMON_TOKEN} PPTX"
    slide.placeholders[1].text = "Teste de conversao de apresentacao do GovGo"

    bullet_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    bullet_slide.shapes.title.text = "Itens"
    text_frame = bullet_slide.placeholders[1].text_frame
    text_frame.text = "Objeto: homologacao MarkItDown"
    paragraph = text_frame.add_paragraph()
    paragraph.text = "Tipo: apresentacao de teste"
    presentation.save(path)
    return {
        "name": "pptx-local",
        "filename": path.name,
        "expected_contains": f"{COMMON_TOKEN} PPTX",
        "expected_absent": ["[Content_Types].xml", "<a:theme", "ppt/slides"],
    }


def _build_pdf_sample() -> Dict[str, str]:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    path = SAMPLES_DIR / "sample_markitdown.pdf"
    pdf = canvas.Canvas(str(path), pagesize=A4)
    pdf.setTitle(f"{COMMON_TOKEN} PDF")
    pdf.drawString(72, 800, f"{COMMON_TOKEN} PDF")
    pdf.drawString(72, 780, "Documento PDF de teste do laboratorio GovGo.")
    pdf.drawString(72, 760, "Objeto: resumo de edital e anexos.")
    pdf.save()
    return {
        "name": "pdf-local",
        "filename": path.name,
        "expected_contains": f"{COMMON_TOKEN} PDF",
    }


def _build_pdf_table_sample() -> Dict[str, str]:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    path = SAMPLES_DIR / "sample_markitdown_table.pdf"
    styles = getSampleStyleSheet()
    story = [
        Paragraph(f"{COMMON_TOKEN} PDF TABLE", styles["Title"]),
        Spacer(1, 12),
    ]
    data = [
        ["Item", "Categoria", "Valor"],
        ["Luva", "EPI", "42"],
        ["Seringa", "Consumo", "128"],
    ]
    table = Table(data)
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.black),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ]
        )
    )
    story.append(table)
    SimpleDocTemplate(str(path), pagesize=A4).build(story)
    return {
        "name": "pdf-table-local",
        "filename": path.name,
        "expected_contains_all": [
            f"{COMMON_TOKEN} PDF TABLE",
            "| Item | Categoria | Valor |",
            "| Luva | EPI | 42 |",
        ],
    }


def ensure_sample_documents() -> List[Dict[str, str]]:
    SAMPLES_DIR.mkdir(parents=True, exist_ok=True)
    cases = _build_text_samples()
    cases.append(_build_pdf_sample())
    cases.append(_build_pdf_table_sample())
    cases.append(_build_docx_sample())
    cases.append(_build_pptx_sample())
    for case in cases:
        case["path"] = str(SAMPLES_DIR / case["filename"])
    return cases